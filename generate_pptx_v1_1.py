"""
Generate PowerPoint Presentation V1.1 for Bob-a-thon Hackathon
Updated with detailed Context Studio schema information
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Travel Refund Uncertainty Estimation System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Refund Estimation with IBM ICA Integration"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Version and details
    details_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    details_frame = details_box.text_frame
    details_frame.text = "Bob ICA Catalysts Hackathon May 2026\nTeam: Abhijit Joshi\nVersion 1.1"
    for para in details_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Hackathon Journey Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Hackathon Journey: 8 Steps to Success"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Project Ideation with Bob AI Assistant"
    
    steps = [
        "2. Building Complete Full-Stack Application",
        "3. Creating Context Graph with Detailed Schema (812 lines)",
        "4. Implementing MCP Server with JSON-RPC 2.0",
        "5. Creating Basic Agent in ICA",
        "6. Building Agent Orchestration with LangGraph",
        "7. Linking Agent to MCP Server",
        "8. Demo Preparation and Testing"
    ]
    
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 3: Step 1 - Project Ideation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 1: Project Ideation with Bob"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Problem Statement"
    
    points = [
        "Travel companies struggle with refund estimation during force majeure",
        "Manual process is inconsistent and time-consuming",
        "Customers need transparent, data-driven estimates",
        "",
        "Solution: AI-Powered System",
        "ML predictions with 95% confidence intervals",
        "Real-time risk assessment",
        "Multi-component booking support",
        "IBM ICA integration for enterprise deployment"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and not point.startswith("Solution") else 0
        p.font.size = Pt(16)
    
    # Slide 4: Step 2 - Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 2: Technology Stack"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Frontend: React 18 + Vite"
    
    points = [
        "7 pages: Dashboard, Bookings, Analytics, Risk Monitor",
        "Interactive charts with Recharts",
        "Responsive design with SCSS",
        "",
        "Backend: FastAPI + Python",
        "11 REST API endpoints",
        "SQLite database with SQLAlchemy ORM",
        "500+ synthetic records for testing",
        "",
        "ML Models: scikit-learn",
        "Random Forest + Gradient Boosting ensemble",
        "95% confidence interval calculation",
        "Uncertainty quantification"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and not point.startswith(("Frontend", "Backend", "ML")) else 0
        p.font.size = Pt(14)
    
    # Slide 5: Step 3 - Context Studio Schema (NEW DETAILED SLIDE)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 3: Context Studio - Detailed Schema"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "context_schema.yaml (812 lines)"
    
    points = [
        "System Architecture: Full-stack components definition",
        "Backend API: 11 endpoints with complete specifications",
        "Database Schema: 6 tables with relationships",
        "Data Models: Request/response schemas with validation",
        "ML Models: Random Forest + Gradient Boosting specs",
        "Business Rules: Refund calculation logic",
        "",
        "refund-api-bulk-import-v2.json (465 lines)",
        "11 REST API tool definitions for Context Forge",
        "Complete JSON Schema validation for each tool",
        "Ready for IBM ICA Agent Orchestration",
        "",
        "Knowledge Graph: 80+ nodes, 179+ relationships",
        "Entities: Bookings, Components, Events, Providers, Policies"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and not point.startswith(("context_schema", "refund-api", "Knowledge")) else 0
        p.font.size = Pt(13)
    
    # Slide 6: Context Schema Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Context Schema: Key Components"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "System Documentation (812 lines)"
    
    points = [
        "Lines 10-18: Architecture (backend, frontend, mcp_server, database, ml_models)",
        "Lines 21-231: Backend API (11 endpoints, 6 database tables)",
        "Lines 234-482: Data Models (BookingCreate, RefundEstimateRequest, etc.)",
        "Lines 523-573: ML Models (Random Forest, Gradient Boosting, Ensemble)",
        "Lines 576-610: Frontend (React pages, routing, API client)",
        "Lines 613-666: MCP Server (11 tools with parameter schemas)",
        "Lines 669-703: Business Rules (refund calculation, risk assessment)",
        "Lines 706-790: Deployment, Security, Performance specifications",
        "",
        "Complete context for AI agents to understand:",
        "API structure, data models, ML capabilities, business logic"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point.startswith("Lines") or point.startswith("API") else 0
        p.font.size = Pt(12)
    
    # Slide 7: Bulk Import Tools
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Bulk Import: 11 API Tools"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "refund-api-bulk-import-v2.json (465 lines)"
    
    points = [
        "v2-create-booking: POST /api/bookings",
        "v2-get-bookings: GET /api/bookings (with pagination)",
        "v2-get-booking-details: GET /api/bookings/{id}",
        "v2-estimate-refund: POST /api/estimate-refund/{id} (ML-powered)",
        "v2-get-refund-estimates: GET /api/estimates/{id}",
        "v2-get-risk-events: GET /api/risk-events (global monitoring)",
        "v2-get-regional-risk: GET /api/risk-events/region/{region}",
        "v2-get-historical-refunds: GET /api/historical-refunds",
        "v2-get-refund-statistics: GET /api/statistics/refund-rates",
        "v2-get-provider-policies: GET /api/providers",
        "v2-get-provider-policy: GET /api/providers/{name}",
        "",
        "Each tool: Complete JSON Schema, headers, authentication, descriptions"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point.startswith("v2-") else 0
        p.font.size = Pt(12)
    
    # Slide 8: Step 4 - MCP Server
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 4: MCP Server Implementation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Model Context Protocol (MCP)"
    
    points = [
        "JSON-RPC 2.0 protocol for AI agent communication",
        "7 tools exposed to IBM ICA agents",
        "Custom implementation in backend/mcp_endpoint_simple.py",
        "",
        "MCP Methods Implemented:",
        "initialize: Handshake with client",
        "tools/list: Return available tools",
        "tools/call: Execute tool and return results",
        "notifications/initialized: Acknowledge initialization",
        "ping: Health check",
        "",
        "Public Access:",
        "ngrok tunnel: https://famished-vertebrae-basil.ngrok-free.dev/mcp",
        "All tools tested and working"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and not point.startswith(("Model", "MCP Methods", "Public")) else 0
        p.font.size = Pt(14)
    
    # Slide 9: Steps 5-7 - ICA Integration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Steps 5-7: IBM ICA Integration"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Step 5: Basic Agent Creation"
    
    points = [
        "Created TravelRefundAdvisor in Agent Studio",
        "Conversational interface configured",
        "",
        "Step 6: Agent Orchestration",
        "LangGraph + ReAct pattern",
        "gpt-5.2 model",
        "All 7 MCP tools integrated",
        "",
        "Step 7: Linking to MCP Server",
        "Registered External MCP Server in ICA Tools",
        "Tool discovery successful (7 tools)",
        "End-to-end integration working",
        "Agent successfully calling tools and retrieving data"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and not point.startswith("Step") else 0
        p.font.size = Pt(14)
    
    # Slide 10: Step 8 - Demo & Testing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 8: Demo Preparation & Testing"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "10 Demo Prompts Tested"
    
    points = [
        "1. Show all travel bookings (53 bookings retrieved)",
        "2. Get booking details for specific ID",
        "3. Estimate refund with ML (95% confidence intervals)",
        "4. Compare severity scenarios",
        "5. Global risk events assessment",
        "6. Regional risk checks (e.g., Ukraine)",
        "7. Historical refund statistics",
        "8. Provider policy comparison",
        "9. Multi-component analysis",
        "10. AI explainability (confidence interval meaning)",
        "",
        "All prompts working successfully!",
        "Terminal logs confirm MCP tool calls",
        "Complete documentation created"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and (point[0].isdigit() or point.startswith(("All", "Terminal", "Complete"))) else 0
        p.font.size = Pt(13)
    
    # Slide 11: Technical Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technical Architecture"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "IBM ICA Platform"
    
    points = [
        "Agent Orchestration (LangGraph + ReAct)",
        "MCP Server (External) - Tool discovery & calling",
        "",
        "Local Development Environment",
        "FastAPI Backend - 11 REST endpoints, MCP endpoint",
        "ML Models - Random Forest + Gradient Boosting",
        "SQLite Database - 6 tables with relationships",
        "React Frontend - 7 pages, interactive charts",
        "",
        "Integration Layer",
        "JSON-RPC 2.0 protocol",
        "ngrok tunnel for public access",
        "HTTPS communication"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and not point.startswith(("IBM", "Local", "Integration")) else 0
        p.font.size = Pt(14)
    
    # Slide 12: Key Achievements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Achievements"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Technical Excellence"
    
    points = [
        "Full-stack application (React + FastAPI)",
        "ML sophistication (ensemble with confidence intervals)",
        "MCP protocol implementation",
        "IBM ICA Agent Orchestration integration",
        "",
        "Innovation",
        "Uncertainty quantification (95% CI, not just predictions)",
        "Force majeure specialized handling",
        "Multi-component granular analysis",
        "AI transparency and explainability",
        "",
        "Business Value",
        "Reduces customer service load (instant estimates)",
        "Improves customer satisfaction (transparent answers)",
        "Real-time risk management",
        "Scalable, production-ready architecture"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and not point.startswith(("Technical", "Innovation", "Business")) else 0
        p.font.size = Pt(13)
    
    # Slide 13: Challenges Overcome
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Challenges Overcome"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. MCP Protocol Implementation"
    
    points = [
        "Problem: Limited documentation, FastMCP compatibility issues",
        "Solution: Custom implementation using JSON-RPC 2.0 directly",
        "",
        "2. ICA Tool Registry",
        "Problem: External MCP Server not appearing in Agent Orchestration",
        "Solution: Discovered separate registries, proper visibility settings",
        "",
        "3. Confidence Interval Calculation",
        "Problem: Single models don't provide uncertainty estimates",
        "Solution: Ensemble approach with prediction variance",
        "",
        "4. Unicode Encoding (Windows)",
        "Problem: Emoji characters causing crashes",
        "Solution: Replaced with ASCII text markers"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point.startswith(("Problem", "Solution")) else 0
        p.font.size = Pt(12)
    
    # Slide 14: Results & Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Results & Impact"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Quantitative Results"
    
    points = [
        "11 API endpoints - All functional",
        "7 MCP tools - Successfully integrated",
        "95% confidence intervals - Uncertainty quantification",
        "53 sample bookings - Realistic test data",
        "500+ historical records - ML training data",
        "100% success rate - All demo prompts working",
        "",
        "Qualitative Impact",
        "Customer Experience: Instant, transparent estimates",
        "Business Efficiency: Automated vs manual process",
        "Risk Management: Real-time global monitoring",
        "AI Transparency: Explainable predictions",
        "Enterprise Ready: IBM ICA integration"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1 if point and not point.startswith(("Quantitative", "Qualitative")) else 0
        p.font.size = Pt(13)
    
    # Slide 15: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Complete System Built in 48 Hours"
    
    points = [
        "",
        "Technical Excellence: Full-stack with ML and enterprise integration",
        "Innovation: Uncertainty quantification and AI transparency",
        "Business Value: Solves real travel industry problem",
        "Production Quality: Professional code, docs, testing",
        "",
        "Integration with IBM ICA through MCP protocol showcases",
        "how modern AI agents can connect to specialized business",
        "applications, creating powerful, transparent, user-friendly solutions.",
        "",
        "This project represents the future of AI-powered",
        "business applications: intelligent, explainable,",
        "and seamlessly integrated."
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)
        p.font.bold = True if point.startswith(("Technical", "Innovation", "Business", "Production")) else False
    
    # Slide 16: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)
    
    # Thank you text
    thank_you_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(60)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = RGBColor(255, 255, 255)
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Bob ICA Catalysts Hackathon May 2026\nTeam: Abhijit Joshi\nVersion 1.1 - Updated with Context Studio Details"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(200, 200, 200)
        para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('Refund_Intelligence_Platform_BoB-a-Thon_V1.1.pptx')
    print('[SUCCESS] PowerPoint V1.1 created: Refund_Intelligence_Platform_BoB-a-Thon_V1.1.pptx')
    print('[INFO] 16 slides with detailed Context Studio schema information')
    print('[INFO] Location: c:/Users/AbhijitJoshi/Uncertainty_Refund/')

if __name__ == '__main__':
    create_presentation()

# Made with Bob
