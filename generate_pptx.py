from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUTPUT_FILE = "Refund_Intelligence_Platform_Demo.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

PALETTE = {
    "bg": RGBColor(255, 250, 246),
    "surface": RGBColor(255, 255, 255),
    "lavender": RGBColor(143, 124, 255),
    "peach": RGBColor(255, 190, 154),
    "mint": RGBColor(173, 232, 214),
    "rose": RGBColor(255, 214, 224),
    "sky": RGBColor(207, 232, 255),
    "text": RGBColor(56, 43, 77),
    "muted": RGBColor(108, 92, 140),
    "line": RGBColor(227, 218, 244),
    "placeholder": RGBColor(244, 239, 255),
}

slides = [
    {
        "title": "Refund Intelligence Platform",
        "subtitle": "AI-assisted refund estimation, risk monitoring, and model comparison",
        "bullets": [
            "Presented by: Abhijit Joshi",
            "Demo duration target: 3-5 minutes",
            "Use this deck for walkthrough, viva, or screen-by-screen demonstration",
        ],
        "speaker": [
            "Introduce the platform as a decision-support solution for refund estimation under uncertainty.",
            "Mention the integration of booking workflows, risk events, historical analytics, and model comparison.",
        ],
        "screenshot": "Cover visual / product branding / optional architecture image",
    },
    {
        "title": "The Problem We Are Solving",
        "bullets": [
            "Travel refund decisions become uncertain during force majeure situations",
            "Customers and operators need visibility into possible refund outcomes",
            "Manual estimation is slow, inconsistent, and difficult to explain",
            "Different calamity types influence refund likelihood differently",
            "Historical outcomes are hard to compare without a unified platform",
        ],
        "speaker": [
            "Position the project as a decision intelligence platform, not just a booking app.",
            "Explain that the goal is to improve confidence and speed in refund-related decisions.",
        ],
        "screenshot": "Optional infographic or title visual",
    },
    {
        "title": "Solution Overview",
        "bullets": [
            "Centralized booking management",
            "Risk event monitoring dashboard",
            "Refund estimation engine",
            "Historical refund analytics",
            "Model comparison workflow",
            "Bell notifications for new bookings and risk events",
        ],
        "speaker": [
            "Summarize the end-to-end workflow supported by the platform.",
            "Explain that users can create data, analyze scenarios, and compare outcomes in one interface.",
        ],
        "screenshot": "Optional collage of app screens",
    },
    {
        "title": "System Architecture",
        "bullets": [
            "Frontend: React + Vite + SCSS",
            "Backend: FastAPI",
            "Database: SQLite (backend/refund_estimation.db)",
            "ML / Estimation: Random Forest, Gradient Boosting, Rule-based logic",
            "Data inputs: bookings, provider policies, historical refunds, risk events",
        ],
        "speaker": [
            "Briefly explain the relationship between frontend, API layer, database, and estimation logic.",
            "Mention that synthetic data is used to bootstrap historical analysis and model behavior.",
        ],
        "screenshot": "Add architecture diagram here",
    },
    {
        "title": "Dashboard Overview",
        "bullets": [
            "Displays summary metrics for bookings, refund trends, and active risks",
            "Provides quick-access navigation to major workflows",
            "Acts as the operational snapshot screen for the platform",
        ],
        "speaker": [
            "Describe the dashboard as the starting point for users and stakeholders.",
            "Highlight the high-level visibility it gives before drilling into detailed screens.",
        ],
        "screenshot": "Insert Dashboard screenshot here",
    },
    {
        "title": "Bookings List",
        "bullets": [
            "Shows all customer bookings in a structured list",
            "Supports review of routes, dates, customer details, and total cost",
            "Provides access to booking details and refund estimation workflows",
        ],
        "speaker": [
            "Explain that bookings are the primary records used for refund analysis.",
            "Mention that this screen organizes customer travel information for operations teams.",
        ],
        "screenshot": "Insert Booking List screenshot here",
    },
    {
        "title": "Create New Booking",
        "bullets": [
            "Captures customer details and travel information",
            "Stores booking component costs such as flight and hotel values",
            "Feeds later refund estimation and comparison workflows",
            "Triggers a bell notification when a new booking is created",
        ],
        "speaker": [
            "Explain that the booking creation form is the main data entry point.",
            "Mention the notification feature for immediate operational visibility.",
        ],
        "screenshot": "Insert Create Booking screenshot here",
    },
    {
        "title": "Booking Details and Refund Estimate",
        "bullets": [
            "Displays complete booking information and cost breakdown",
            "Generates refund estimate for the selected booking",
            "Shows refund amount, refund percentage, confidence range, and risk score",
            "Presents best-case, likely, and worst-case views",
        ],
        "speaker": [
            "Use this slide to explain how a single booking can be evaluated in detail.",
            "Emphasize scenario visibility and explainability of the result.",
        ],
        "screenshot": "Insert Booking Details screenshot here",
    },
    {
        "title": "Refund Model Comparator",
        "bullets": [
            "Select an existing booking",
            "Choose estimation model: Auto, Random Forest, Gradient Boosting, or Rule-based",
            "Choose calamity type and severity",
            "Compare refund outcomes for the selected scenario",
        ],
        "speaker": [
            "Highlight this as a key differentiator of the project.",
            "Explain how comparator support allows scenario testing and model selection flexibility.",
        ],
        "screenshot": "Insert Refund Comparator screenshot here",
    },
    {
        "title": "Risk Monitor and Event Entry",
        "bullets": [
            "Tracks disruptions such as war, natural disaster, pandemic, or strikes",
            "Displays severity, region, date range, and status of risk events",
            "Allows users to add new risk events through a form",
            "New risk event creation triggers a bell notification",
        ],
        "speaker": [
            "Explain how real-world disruptions are modeled as risk events in the platform.",
            "Mention that these events directly support risk-aware refund reasoning.",
        ],
        "screenshot": "Insert Risk Monitor screenshot here",
    },
    {
        "title": "Bell Notifications",
        "bullets": [
            "Notification badge appears on the header bell icon",
            "Triggered when a new booking is created",
            "Triggered when a new risk event is created",
            "Dropdown shows recent notifications with clear action",
        ],
        "speaker": [
            "Describe notifications as a lightweight operational awareness feature.",
            "Mention that current notifications are session-based and can be extended later.",
        ],
        "screenshot": "Insert Notification dropdown screenshot here",
    },
    {
        "title": "Historical Refund Analytics",
        "bullets": [
            "Visualizes historical refund patterns and averages",
            "Supports filtering by component or event type",
            "Helps interpret model behavior using historical trends",
            "Strengthens explainability beyond a single prediction",
        ],
        "speaker": [
            "Present analytics as the explanatory layer of the platform.",
            "Explain that it helps users understand past outcomes, not just current estimates.",
        ],
        "screenshot": "Insert Analytics screenshot here",
    },
    {
        "title": "How Refund Estimation Works",
        "bullets": [
            "Booking components are evaluated individually",
            "Risk score is derived from active events or selected calamity and severity",
            "The engine supports Auto Ensemble, Random Forest, Gradient Boosting, and Rule-based fallback",
            "Final refund result aggregates component-level reasoning into one estimate",
        ],
        "speaker": [
            "Summarize the backend decision process in simple language.",
            "Note that rule-based logic provides resilience when ML comparison is not preferred.",
        ],
        "screenshot": "Optional flow diagram here",
    },
    {
        "title": "Suggested 3-5 Minute Demo Flow",
        "bullets": [
            "Start with Dashboard",
            "Open Bookings List",
            "Create a new booking and show notification",
            "Open Booking Details and generate estimate",
            "Open Comparator and compare models",
            "Add a new event in Risk Monitor and show notification",
            "End with Analytics for historical interpretation",
        ],
        "speaker": [
            "This slide can guide your live or recorded narration.",
            "Keep the flow concise and focused on user value and system capabilities.",
        ],
        "screenshot": "Optional end-to-end workflow image",
    },
    {
        "title": "Business Value",
        "bullets": [
            "Faster refund decision support",
            "Better visibility into uncertainty and disruptions",
            "Model-based comparison for scenario analysis",
            "Centralized platform for bookings, events, and analytics",
            "Improved communication for customers and operators",
        ],
        "speaker": [
            "Explain why the platform matters from an operational and analytical perspective.",
            "Position the project as both practical and extensible.",
        ],
        "screenshot": "Optional impact graphic here",
    },
    {
        "title": "Future Enhancements",
        "bullets": [
            "Persistent notifications",
            "External live risk feeds",
            "Provider-specific refund policy intelligence",
            "Model evaluation metrics and benchmarking",
            "Authentication and role-based access",
            "Exportable reports and comparison snapshots",
        ],
        "speaker": [
            "Show that the platform can evolve into a production-grade decision-support system.",
            "Mention the most realistic next steps for future scope.",
        ],
        "screenshot": "Optional roadmap visual here",
    },
    {
        "title": "Thank You",
        "subtitle": "Questions? Demo complete.",
        "bullets": [
            "Refund Intelligence Platform",
            "AI-assisted refund estimation under uncertainty",
        ],
        "speaker": [
            "Close the presentation confidently and invite questions.",
        ],
        "screenshot": "Optional closing visual",
    },
]

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PALETTE["bg"]
    bg.line.fill.background()

    top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = PALETTE["lavender"]
    top_band.line.fill.background()

    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.6), Inches(0.35), Inches(1.15), Inches(1.15))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = PALETTE["rose"]
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.9), Inches(0.9), Inches(0.75), Inches(0.75))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = PALETTE["mint"]
    circle2.line.fill.background()

def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(8.4), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = PALETTE["text"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.28), Inches(8.0), Inches(0.5))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        r2.font.color.rgb = PALETTE["muted"]

def add_content_box(slide, bullets):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.85), Inches(6.0), Inches(4.85))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALETTE["surface"]
    shape.line.color.rgb = PALETTE["line"]

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12)
    tf.clear()

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18 if idx == 0 and len(bullets) <= 3 else 16)
        p.font.color.rgb = PALETTE["text"]
        p.space_after = Pt(10)
        p.bullet = True

def add_speaker_box(slide, speaker_lines):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.42))
    box.fill.solid()
    box.fill.fore_color.rgb = PALETTE["sky"]
    box.line.color.rgb = PALETTE["line"]

    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Speaker cue: " + " ".join(speaker_lines)
    run.font.size = Pt(10.5)
    run.font.color.rgb = PALETTE["text"]

def add_placeholder(slide, caption):
    x = Inches(6.9)
    y = Inches(1.85)
    w = Inches(5.85)
    h = Inches(4.85)

    ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    ph.fill.solid()
    ph.fill.fore_color.rgb = PALETTE["placeholder"]
    ph.line.color.rgb = PALETTE["line"]

    inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.28), y + Inches(0.32), w - Inches(0.56), h - Inches(1.0))
    inner.fill.solid()
    inner.fill.fore_color.rgb = PALETTE["surface"]
    inner.line.color.rgb = PALETTE["line"]

    label = slide.shapes.add_textbox(x + Inches(0.35), y + Inches(0.45), w - Inches(0.7), Inches(0.5))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Screenshot Placeholder"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = PALETTE["lavender"]

    cap = slide.shapes.add_textbox(x + Inches(0.45), y + Inches(3.95), w - Inches(0.9), Inches(0.55))
    tf2 = cap.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = caption
    r2.font.size = Pt(12)
    r2.font.color.rgb = PALETTE["muted"]

for item in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, item["title"], item.get("subtitle"))
    add_content_box(slide, item.get("bullets", []))
    add_placeholder(slide, item.get("screenshot", "Add relevant screenshot here"))
    add_speaker_box(slide, item.get("speaker", []))

prs.save(OUTPUT_FILE)
print(f"Created {OUTPUT_FILE} with {len(slides)} slides.")

# Made with Bob
